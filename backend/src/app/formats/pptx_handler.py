from pathlib import Path

from pptx import Presentation

from app.formats.base import FormatHandler, register
from app.models.job import TextSegment


@register([".pptx"])
class PptxHandler(FormatHandler):
    def extract_texts(self, file_path: Path) -> list[TextSegment]:
        prs = Presentation(str(file_path))
        segments = []
        for s_idx, slide in enumerate(prs.slides):
            slide_title = ""
            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text

            for sh_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for r_idx, run in enumerate(paragraph.runs):
                        if run.text.strip():
                            segments.append(TextSegment(
                                id=f"s{s_idx}:sh{sh_idx}:p{p_idx}:r{r_idx}",
                                text=run.text,
                                context=f"Slide: {slide_title}" if slide_title else "",
                                metadata={
                                    "slide": s_idx,
                                    "shape": sh_idx,
                                    "paragraph": p_idx,
                                    "run": r_idx,
                                },
                            ))
        return segments

    def replace_texts(self, file_path: Path, translated: list[TextSegment], output_path: Path) -> None:
        prs = Presentation(str(file_path))
        lookup = {seg.id: seg.text for seg in translated}

        for s_idx, slide in enumerate(prs.slides):
            for sh_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for r_idx, run in enumerate(paragraph.runs):
                        key = f"s{s_idx}:sh{sh_idx}:p{p_idx}:r{r_idx}"
                        if key in lookup:
                            run.text = lookup[key]

        prs.save(str(output_path))
